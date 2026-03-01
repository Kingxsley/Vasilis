"""
Executive Training PowerPoint Generator
Generates detailed 30-50 slide presentations for security awareness training modules
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io
from typing import Dict, List, Optional
import logging

logger = logging.getLogger(__name__)

# Color scheme - Professional dark theme
COLORS = {
    "primary": RGBColor(212, 168, 54),      # Gold #D4A836
    "background": RGBColor(13, 17, 23),     # Dark #0D1117
    "text": RGBColor(232, 221, 181),        # Light text #E8DDB5
    "accent": RGBColor(41, 121, 255),       # Blue #2979FF
    "danger": RGBColor(255, 59, 48),        # Red #FF3B30
    "success": RGBColor(0, 230, 118),       # Green #00E676
    "warning": RGBColor(255, 179, 0),       # Yellow #FFB300
    "white": RGBColor(255, 255, 255),
    "dark_bg": RGBColor(22, 27, 34),        # Card bg #161B22
}


def add_slide(prs, layout_index=6):
    """Add a blank slide to the presentation"""
    layout = prs.slide_layouts[layout_index]  # 6 = blank
    return prs.slides.add_slide(layout)


def set_slide_background(slide, color=COLORS["background"]):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=0.5, left=0.5, width=9, height=1, font_size=44, bold=True, color=COLORS["primary"]):
    """Add a title to the slide"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return shape


def add_content(slide, text, top=2, left=0.5, width=9, height=5, font_size=24, color=COLORS["text"]):
    """Add content text to the slide"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Arial"
    return shape


def add_bullet_list(slide, items: List[str], top=2, left=0.5, width=9, height=5, font_size=20, color=COLORS["text"]):
    """Add a bullet list to the slide"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(12)
    return shape


def add_subtitle(slide, text, top=1.2, left=0.5, width=9, height=0.8, font_size=24, color=COLORS["text"]):
    """Add subtitle to the slide"""
    return add_title(slide, text, top, left, width, height, font_size, False, color)


def add_icon_card(slide, icon_text, title, description, left, top, width=4, height=2):
    """Add an icon card for statistics or features"""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["dark_bg"]
    shape.line.fill.background()
    
    # Icon/Number
    icon_shape = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(0.6))
    tf = icon_shape.text_frame
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    
    # Title
    title_shape = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.8), Inches(width - 0.4), Inches(0.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    
    # Description
    desc_shape = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.3), Inches(width - 0.4), Inches(0.6))
    tf = desc_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)


def add_two_column(slide, left_title, left_items, right_title, right_items, top=2):
    """Add two column layout"""
    # Left column title
    add_title(slide, left_title, top, 0.5, 4, 0.6, 24, True, COLORS["primary"])
    add_bullet_list(slide, left_items, top + 0.7, 0.5, 4, 3, 16, COLORS["text"])
    
    # Right column title
    add_title(slide, right_title, top, 5, 4, 0.6, 24, True, COLORS["primary"])
    add_bullet_list(slide, right_items, top + 0.7, 5, 4, 3, 16, COLORS["text"])


def add_warning_box(slide, text, top=5, left=0.5, width=9, height=1):
    """Add a warning/alert box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 59, 48)  # Red
    shape.fill.fore_color.brightness = 0.8  # Lighter
    shape.line.fill.background()
    
    text_shape = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.2), Inches(width - 0.6), Inches(height - 0.4))
    tf = text_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"⚠️ {text}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS["danger"]


# ================== Module-Specific Content ==================

PHISHING_CONTENT = {
    "title": "Email Phishing Awareness",
    "subtitle": "Protecting Against Social Engineering Attacks",
    "slides": [
        {
            "title": "What is Phishing?",
            "content": "Phishing is a type of social engineering attack where attackers impersonate legitimate entities to trick victims into revealing sensitive information.",
            "bullets": [
                "Fraudulent emails designed to steal credentials",
                "Impersonation of trusted brands and services",
                "Often creates sense of urgency",
                "Can lead to data breaches and financial loss"
            ]
        },
        {
            "title": "Phishing Statistics",
            "type": "stats",
            "stats": [
                ("91%", "Cyber Attacks", "Start with phishing email"),
                ("$4.91M", "Average Cost", "Of a data breach"),
            ]
        },
        {
            "title": "Common Phishing Techniques",
            "bullets": [
                "Urgency and fear tactics",
                "Spoofed sender addresses",
                "Lookalike domains",
                "Embedded malicious links",
                "Fake login pages",
                "Attachment-based attacks"
            ]
        },
        {
            "title": "Red Flags: Email Headers",
            "bullets": [
                "Suspicious sender email address",
                "Mismatched 'From' and 'Reply-To'",
                "Generic greetings ('Dear Customer')",
                "Unexpected attachments",
                "External sender warnings"
            ]
        },
        {
            "title": "Red Flags: Email Content",
            "bullets": [
                "Spelling and grammar errors",
                "Threatening or urgent language",
                "Requests for sensitive information",
                "Too-good-to-be-true offers",
                "Suspicious URLs when hovering"
            ]
        },
        {
            "title": "Types of Phishing",
            "type": "two_column",
            "left_title": "Common Types",
            "left_items": ["Spear Phishing", "Whaling", "Clone Phishing", "Vishing (Voice)"],
            "right_title": "Advanced Types",
            "right_items": ["Smishing (SMS)", "Business Email Compromise", "Angler Phishing", "Pharming"]
        },
        {
            "title": "Spear Phishing Deep Dive",
            "content": "Targeted attacks using personal information about the victim",
            "bullets": [
                "Researched using social media",
                "Highly personalized content",
                "Often targets executives",
                "More convincing and dangerous"
            ]
        },
        {
            "title": "Business Email Compromise (BEC)",
            "bullets": [
                "Attackers impersonate executives",
                "Request wire transfers or sensitive data",
                "Use spoofed or compromised accounts",
                "Cost businesses billions annually",
                "Often bypasses technical controls"
            ]
        },
        {
            "title": "How to Verify Suspicious Emails",
            "bullets": [
                "Check sender's actual email address",
                "Hover over links before clicking",
                "Contact sender through known channels",
                "Look for HTTPS and valid certificates",
                "Report suspicious emails to IT"
            ]
        },
        {
            "title": "Safe Browsing Practices",
            "bullets": [
                "Never click links in suspicious emails",
                "Type URLs directly in browser",
                "Verify website security (HTTPS)",
                "Use bookmarks for sensitive sites",
                "Keep browser updated"
            ]
        },
        {
            "title": "Protecting Your Credentials",
            "bullets": [
                "Use unique passwords for each site",
                "Enable multi-factor authentication",
                "Never share passwords via email",
                "Use a password manager",
                "Change passwords after breaches"
            ]
        },
        {
            "title": "What To Do If You Clicked",
            "bullets": [
                "Disconnect from network immediately",
                "Report to IT security team",
                "Change affected passwords",
                "Monitor accounts for suspicious activity",
                "Document what happened"
            ],
            "warning": "Time is critical - report immediately!"
        },
        {
            "title": "Reporting Phishing",
            "bullets": [
                "Forward suspicious emails to IT",
                "Use phishing report button if available",
                "Don't forward to colleagues",
                "Delete after reporting",
                "Report to external authorities if needed"
            ]
        },
        {
            "title": "Case Study: Real Attack",
            "content": "A major corporation lost $100 million to BEC attack",
            "bullets": [
                "Attackers posed as supplier",
                "Changed bank account details",
                "Bypassed verification procedures",
                "Funds transferred overseas",
                "Recovery was partial"
            ]
        },
        {
            "title": "Best Practices Summary",
            "bullets": [
                "Think before you click",
                "Verify unexpected requests",
                "Use multi-factor authentication",
                "Keep systems updated",
                "Report suspicious activity",
                "Stay informed about new threats"
            ]
        }
    ]
}

SOCIAL_ENGINEERING_CONTENT = {
    "title": "Social Engineering Defense",
    "subtitle": "Understanding and Preventing Human-Targeted Attacks",
    "slides": [
        {
            "title": "What is Social Engineering?",
            "content": "Psychological manipulation techniques used to trick people into revealing confidential information or performing actions that compromise security.",
            "bullets": [
                "Exploits human psychology, not technology",
                "Targets trust, fear, and urgency",
                "Can bypass technical security measures",
                "Affects organizations of all sizes"
            ]
        },
        {
            "title": "Why Social Engineering Works",
            "bullets": [
                "Humans are the weakest link",
                "People want to be helpful",
                "Authority figures are trusted",
                "Time pressure reduces critical thinking",
                "Social norms encourage compliance"
            ]
        },
        {
            "title": "Psychological Principles",
            "type": "two_column",
            "left_title": "Manipulation Tactics",
            "left_items": ["Authority", "Scarcity", "Social Proof", "Reciprocity"],
            "right_title": "Emotional Triggers",
            "right_items": ["Fear", "Curiosity", "Greed", "Helpfulness"]
        },
        {
            "title": "Types of Social Engineering",
            "bullets": [
                "Pretexting - creating false scenarios",
                "Baiting - offering something enticing",
                "Quid Pro Quo - exchange of services",
                "Tailgating - physical following",
                "Impersonation - posing as authority"
            ]
        },
        {
            "title": "Pretexting Attacks",
            "content": "Creating a fabricated scenario to engage victims",
            "bullets": [
                "Impersonating IT support",
                "Posing as vendors or contractors",
                "Claiming to be from HR or management",
                "Using fake emergency situations"
            ]
        },
        {
            "title": "Baiting Attacks",
            "bullets": [
                "USB drives left in parking lots",
                "Free software downloads",
                "Fake prize notifications",
                "Infected media devices",
                "Malicious QR codes"
            ],
            "warning": "Never plug in unknown USB devices!"
        },
        {
            "title": "Physical Social Engineering",
            "bullets": [
                "Tailgating through secure doors",
                "Impersonating delivery personnel",
                "Dumpster diving for information",
                "Shoulder surfing passwords",
                "Installing rogue devices"
            ]
        },
        {
            "title": "Phone-Based Attacks (Vishing)",
            "bullets": [
                "Caller ID spoofing",
                "Tech support scams",
                "Bank fraud calls",
                "Government impersonation",
                "Emergency family scams"
            ]
        },
        {
            "title": "Protecting Against Social Engineering",
            "bullets": [
                "Verify identity before sharing information",
                "Follow established procedures",
                "Don't be rushed into decisions",
                "Report suspicious requests",
                "Trust your instincts"
            ]
        },
        {
            "title": "Verification Best Practices",
            "bullets": [
                "Use callback verification",
                "Confirm through official channels",
                "Check employee directories",
                "Require proper authorization",
                "Document unusual requests"
            ]
        },
        {
            "title": "Building a Security Culture",
            "bullets": [
                "Regular security awareness training",
                "Clear reporting procedures",
                "No-blame reporting policy",
                "Regular simulated attacks",
                "Reward security-conscious behavior"
            ]
        },
        {
            "title": "Case Study: Corporate Breach",
            "content": "A Fortune 500 company was breached through a phone call",
            "bullets": [
                "Attacker posed as IT support",
                "Convinced employee to reset password",
                "Gained access to financial systems",
                "Exfiltrated customer data",
                "Cost: $50 million in damages"
            ]
        },
        {
            "title": "Key Takeaways",
            "bullets": [
                "Anyone can be a target",
                "Verify before trusting",
                "Follow security procedures",
                "Report suspicious activity",
                "Stay vigilant at all times"
            ]
        }
    ]
}

PASSWORD_SECURITY_CONTENT = {
    "title": "Password Security",
    "subtitle": "Creating and Managing Strong Credentials",
    "slides": [
        {
            "title": "Why Password Security Matters",
            "content": "Passwords are the first line of defense against unauthorized access to your accounts and sensitive data.",
            "bullets": [
                "81% of breaches involve weak passwords",
                "Average person has 100+ online accounts",
                "Password reuse is extremely common",
                "Stolen credentials are traded on dark web"
            ]
        },
        {
            "title": "Common Password Mistakes",
            "bullets": [
                "Using simple, guessable passwords",
                "Reusing passwords across sites",
                "Writing passwords on sticky notes",
                "Sharing passwords with others",
                "Not changing default passwords"
            ]
        },
        {
            "title": "Password Cracking Methods",
            "type": "two_column",
            "left_title": "Automated Attacks",
            "left_items": ["Brute Force", "Dictionary Attack", "Rainbow Tables", "Credential Stuffing"],
            "right_title": "Social Attacks",
            "right_items": ["Phishing", "Shoulder Surfing", "Social Engineering", "Dumpster Diving"]
        },
        {
            "title": "Creating Strong Passwords",
            "bullets": [
                "Minimum 12 characters (16+ recommended)",
                "Mix uppercase, lowercase, numbers, symbols",
                "Avoid personal information",
                "Don't use dictionary words",
                "Use passphrases instead"
            ]
        },
        {
            "title": "Passphrase Example",
            "content": "A passphrase is easier to remember and harder to crack",
            "bullets": [
                "Bad: 'Password123!'",
                "Good: 'MyDogLoves2RunInThe!Park'",
                "Better: 'Correct-Horse-Battery-Staple-99'",
                "Use unrelated words",
                "Add numbers and symbols"
            ]
        },
        {
            "title": "Password Managers",
            "bullets": [
                "Generate strong, unique passwords",
                "Securely store all credentials",
                "Auto-fill login forms",
                "Sync across devices",
                "Only remember one master password"
            ]
        },
        {
            "title": "Multi-Factor Authentication",
            "content": "Adding a second layer of security beyond passwords",
            "bullets": [
                "Something you know (password)",
                "Something you have (phone, key)",
                "Something you are (biometrics)",
                "Blocks 99.9% of automated attacks"
            ]
        },
        {
            "title": "Types of MFA",
            "type": "two_column",
            "left_title": "Good Options",
            "left_items": ["Authenticator Apps", "Hardware Keys", "Biometrics", "Push Notifications"],
            "right_title": "Less Secure",
            "right_items": ["SMS Codes", "Email Codes", "Security Questions", "Backup Codes"]
        },
        {
            "title": "Password Hygiene",
            "bullets": [
                "Never share passwords",
                "Change passwords after breaches",
                "Use unique password per account",
                "Enable MFA everywhere possible",
                "Regularly audit your accounts"
            ]
        },
        {
            "title": "Checking for Breaches",
            "bullets": [
                "Use haveibeenpwned.com",
                "Enable breach notifications",
                "Monitor dark web exposure",
                "Change compromised passwords immediately",
                "Review connected accounts"
            ],
            "warning": "If your password appears in a breach, change it immediately!"
        },
        {
            "title": "Corporate Password Policy",
            "bullets": [
                "Minimum length requirements",
                "Complexity requirements",
                "Regular password rotation",
                "Account lockout policies",
                "Single sign-on (SSO) usage"
            ]
        },
        {
            "title": "Key Takeaways",
            "bullets": [
                "Use long, unique passwords",
                "Enable MFA on all accounts",
                "Use a password manager",
                "Never reuse passwords",
                "Report suspected compromises"
            ]
        }
    ]
}

DATA_PROTECTION_CONTENT = {
    "title": "Data Protection & Privacy",
    "subtitle": "Safeguarding Sensitive Information",
    "slides": [
        {
            "title": "Why Data Protection Matters",
            "content": "Protecting sensitive data is critical for maintaining trust, compliance, and business continuity.",
            "bullets": [
                "Data breaches cost millions in damages",
                "Regulatory fines can be severe",
                "Reputation damage is long-lasting",
                "Personal liability for negligence"
            ]
        },
        {
            "title": "Types of Sensitive Data",
            "type": "two_column",
            "left_title": "Personal Data",
            "left_items": ["Names & Addresses", "Social Security Numbers", "Financial Information", "Health Records"],
            "right_title": "Business Data",
            "right_items": ["Trade Secrets", "Client Information", "Financial Reports", "Strategic Plans"]
        },
        {
            "title": "Data Classification",
            "bullets": [
                "Public - Can be freely shared",
                "Internal - For employees only",
                "Confidential - Need-to-know basis",
                "Restricted - Highest protection",
                "Always classify before handling"
            ]
        },
        {
            "title": "Regulatory Requirements",
            "bullets": [
                "GDPR - European data protection",
                "CCPA - California consumer privacy",
                "HIPAA - Healthcare information",
                "PCI-DSS - Payment card data",
                "SOX - Financial reporting"
            ],
            "warning": "Non-compliance can result in significant fines!"
        },
        {
            "title": "Data Handling Best Practices",
            "bullets": [
                "Minimize data collection",
                "Encrypt sensitive data",
                "Use secure transmission",
                "Limit access to need-to-know",
                "Properly dispose of data"
            ]
        },
        {
            "title": "Secure Data Storage",
            "bullets": [
                "Use encrypted drives",
                "Implement access controls",
                "Regular backups",
                "Secure cloud storage",
                "Physical security measures"
            ]
        },
        {
            "title": "Data Transmission Security",
            "bullets": [
                "Use encrypted email",
                "Verify recipient before sending",
                "Use secure file sharing",
                "Avoid public WiFi for sensitive data",
                "Use VPN when remote"
            ]
        },
        {
            "title": "Clean Desk Policy",
            "bullets": [
                "Lock computer when away",
                "Secure documents in drawers",
                "Shred sensitive papers",
                "Clear whiteboards after meetings",
                "Don't leave screens visible"
            ]
        },
        {
            "title": "Mobile Device Security",
            "bullets": [
                "Enable device encryption",
                "Use strong screen locks",
                "Enable remote wipe",
                "Keep devices updated",
                "Avoid public charging stations"
            ]
        },
        {
            "title": "Data Breach Response",
            "bullets": [
                "Report immediately to IT",
                "Document what happened",
                "Preserve evidence",
                "Don't try to cover up",
                "Follow incident procedures"
            ]
        },
        {
            "title": "Key Takeaways",
            "bullets": [
                "Know your data classification",
                "Encrypt sensitive information",
                "Follow secure handling procedures",
                "Report incidents immediately",
                "Stay compliant with regulations"
            ]
        }
    ]
}

# Module content mapping
MODULE_CONTENT = {
    "phishing": PHISHING_CONTENT,
    "email_phishing": PHISHING_CONTENT,
    "social_engineering": SOCIAL_ENGINEERING_CONTENT,
    "password_security": PASSWORD_SECURITY_CONTENT,
    "password": PASSWORD_SECURITY_CONTENT,
    "data_protection": DATA_PROTECTION_CONTENT,
    "data_handling": DATA_PROTECTION_CONTENT,
    "privacy": DATA_PROTECTION_CONTENT,
}


def generate_module_presentation(module_name: str, module_content: Optional[Dict] = None) -> bytes:
    """
    Generate a PowerPoint presentation for a training module.
    
    Args:
        module_name: Name or type of the module
        module_content: Optional custom content dictionary
        
    Returns:
        bytes: PowerPoint file as bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Determine content to use
    content = module_content
    if not content:
        # Try to match module name to predefined content
        module_key = module_name.lower().replace(" ", "_").replace("-", "_")
        for key in MODULE_CONTENT:
            if key in module_key or module_key in key:
                content = MODULE_CONTENT[key]
                break
    
    if not content:
        # Generate generic presentation
        content = {
            "title": f"{module_name} Training",
            "subtitle": "Security Awareness Module",
            "slides": [
                {
                    "title": "Introduction",
                    "content": f"Welcome to {module_name} training",
                    "bullets": [
                        "Understanding key concepts",
                        "Best practices",
                        "Recognizing threats",
                        "Response procedures"
                    ]
                },
                {
                    "title": "Key Concepts",
                    "bullets": [
                        "Security is everyone's responsibility",
                        "Stay vigilant at all times",
                        "Report suspicious activity",
                        "Follow established procedures"
                    ]
                },
                {
                    "title": "Best Practices",
                    "bullets": [
                        "Think before you click",
                        "Verify unusual requests",
                        "Keep systems updated",
                        "Use strong authentication"
                    ]
                },
                {
                    "title": "Key Takeaways",
                    "bullets": [
                        "Security starts with you",
                        "When in doubt, verify",
                        "Report incidents immediately",
                        "Stay informed about threats"
                    ]
                }
            ]
        }
    
    # ===== Title Slide =====
    slide = add_slide(prs)
    set_slide_background(slide)
    add_title(slide, content["title"], 2.5, 0.5, 9, 1.5, 54, True, COLORS["primary"])
    add_subtitle(slide, content.get("subtitle", "Security Awareness Training"), 4, 0.5, 9, 1, 28, COLORS["text"])
    add_content(slide, "Executive Training Presentation", 5.5, 0.5, 9, 1, 18, RGBColor(128, 128, 128))
    
    # ===== Agenda Slide =====
    slide = add_slide(prs)
    set_slide_background(slide)
    add_title(slide, "Agenda", 0.5)
    agenda_items = [s.get("title", "Section") for s in content["slides"][:10]]
    add_bullet_list(slide, agenda_items, 1.8)
    
    # ===== Content Slides =====
    for slide_data in content["slides"]:
        slide = add_slide(prs)
        set_slide_background(slide)
        
        slide_type = slide_data.get("type", "standard")
        
        if slide_type == "stats":
            add_title(slide, slide_data["title"], 0.5)
            stats = slide_data.get("stats", [])
            for i, (number, label, desc) in enumerate(stats):
                add_icon_card(slide, number, label, desc, 0.5 + (i * 4.5), 2)
                
        elif slide_type == "two_column":
            add_title(slide, slide_data["title"], 0.5)
            add_two_column(
                slide,
                slide_data.get("left_title", ""),
                slide_data.get("left_items", []),
                slide_data.get("right_title", ""),
                slide_data.get("right_items", [])
            )
        else:
            add_title(slide, slide_data["title"], 0.5)
            if slide_data.get("content"):
                add_content(slide, slide_data["content"], 1.5, 0.5, 9, 1, 22, COLORS["text"])
                add_bullet_list(slide, slide_data.get("bullets", []), 2.5, 0.5)
            else:
                add_bullet_list(slide, slide_data.get("bullets", []), 1.8, 0.5)
            
            if slide_data.get("warning"):
                add_warning_box(slide, slide_data["warning"], 5.5)
    
    # ===== Summary Slide =====
    slide = add_slide(prs)
    set_slide_background(slide)
    add_title(slide, "Summary", 0.5)
    add_bullet_list(slide, [
        "Security is a shared responsibility",
        "Stay vigilant and report suspicious activity",
        "Follow established procedures and policies",
        "Keep learning and stay informed",
        "When in doubt, ask IT Security"
    ], 1.8)
    
    # ===== Q&A Slide =====
    slide = add_slide(prs)
    set_slide_background(slide)
    add_title(slide, "Questions?", 2.5, 0.5, 9, 1.5, 54)
    add_subtitle(slide, "Thank you for your attention", 4, 0.5, 9, 1, 28, COLORS["text"])
    add_content(slide, "Contact IT Security for any concerns", 5.5, 0.5, 9, 1, 18, RGBColor(128, 128, 128))
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def get_available_modules() -> List[str]:
    """Get list of available predefined modules"""
    return list(MODULE_CONTENT.keys())
